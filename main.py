from fastapi import FastAPI, UploadFile, File, Form
from fastapi.responses import FileResponse, JSONResponse
from fastapi.middleware.cors import CORSMiddleware
from openai import OpenAI

import os
import tempfile
from pptx import Presentation
from PyPDF2 import PdfReader
import docx
import pandas as pd
from docx import Document

# FastAPI 앱
app = FastAPI()

# 🔹 CORS 허용 (GitHub Pages 주소만)
app.add_middleware(
    CORSMiddleware,
    allow_origins=["https://bh-1113.github.io"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# OpenAI 클라이언트
client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))

# ======================================
# 0) 루트 엔드포인트 (health check)
# ======================================
@app.get("/")
def root():
    return {
        "message": "AI Report API is running 🚀",
        "endpoints": ["/make_ppt", "/upload_summary"]
    }

# ======================================
# 1) 보고서 자동 생성 (report.html → /make_ppt)
# ======================================
sections = ["개요", "필요성", "활용 사례", "장점과 한계", "미래 전망"]

def generate_text(topic, section):
    prompt = f"{topic}에 대해 '{section}' 파트의 발표 슬라이드 내용을 글머리표 3~4개로 작성해줘."
    response = client.chat.completions.create(
        model="gpt-4o-mini",
        messages=[{"role": "user", "content": prompt}]
    )
    return response.choices[0].message.content

@app.get("/make_ppt")
def make_ppt(topic: str):
    prs = Presentation()

    # 표지
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = f"{topic} 보고서"
    slide.placeholders[1].text = "자동 생성된 AI 보고서"

    all_text = ""  # 웹페이지용 미리보기 텍스트

    # 본문
    for section in sections:
        text = generate_text(topic, section)
        all_text += f"{section}\n{text}\n\n"

        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = section
        slide.placeholders[1].text = text

    # 파일 저장
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
    prs.save(tmp.name)

    return JSONResponse({
        "preview": all_text.strip(),
        "ppt_url": f"/download_ppt/{os.path.basename(tmp.name)}"
    })

@app.get("/download_ppt/{filename}")
def download_ppt(filename: str):
    path = os.path.join(tempfile.gettempdir(), filename)
    if os.path.exists(path):
        return FileResponse(
            path,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            filename=filename
        )
    return JSONResponse({"error": "파일이 존재하지 않습니다."})

# ======================================
# 2) 문서 요약 (summary.html → /upload_summary)
# ======================================
def extract_text(file: UploadFile):
    ext = file.filename.split(".")[-1].lower()
    text = ""

    with tempfile.NamedTemporaryFile(delete=False) as tmp:
        tmp.write(file.file.read())
        tmp_path = tmp.name

    if ext == "pdf":
        reader = PdfReader(tmp_path)
        for page in reader.pages:
            text += page.extract_text() + "\n"

    elif ext == "docx":
        doc = docx.Document(tmp_path)
        for para in doc.paragraphs:
            text += para.text + "\n"

    elif ext == "pptx":
        prs = Presentation(tmp_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"

    elif ext in ["xlsx", "xls"]:
        df = pd.read_excel(tmp_path)
        text = df.to_string()

    elif ext == "txt":
        with open(tmp_path, "r", encoding="utf-8") as f:
            text = f.read()

    else:
        text = "지원하지 않는 파일 형식입니다."

    os.remove(tmp_path)
    return text.strip()

def gpt_summarize(text: str) -> str:
    prompt = f"""
    다음 문서를 간결하게 요약해 주세요. 
    핵심 내용만 남기고, 불필요한 부분은 제거하세요.

    문서:
    {text[:4000]}
    """
    response = client.chat.completions.create(
        model="gpt-4o-mini",
        messages=[{"role": "user", "content": prompt}],
        temperature=0.3,
    )
    return response.choices[0].message.content.strip()

def save_as_docx(summary: str, filename: str):
    doc = Document()
    doc.add_heading("문서 요약", level=1)
    doc.add_paragraph(summary)
    doc.save(filename)

def save_as_pptx(summary: str, filename: str):
    prs = Presentation()
    slide_layout = prs.slide_layouts[1]

    # 제목 슬라이드
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "문서 요약"
    slide.placeholders[1].text = "AI가 자동 생성한 요약 자료"

    # 본문 슬라이드
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "요약 내용"
    content.text = summary

    prs.save(filename)

@app.post("/upload_summary")
async def upload_summary(file: UploadFile = File(...), export: str = Form("json")):
    text = extract_text(file)
    summary = gpt_summarize(text)

    if export == "json":
        return JSONResponse({"summary": summary})
    elif export == "docx":
        tmp_path = "summary.docx"
        save_as_docx(summary, tmp_path)
        return FileResponse(tmp_path, media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document", filename="summary.docx")
    elif export == "pptx":
        tmp_path = "summary.pptx"
        save_as_pptx(summary, tmp_path)
        return FileResponse(tmp_path, media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation", filename="summary.pptx")
    else:
        return JSONResponse({"error": "지원하지 않는 export 형식입니다."})
